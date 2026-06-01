#!/usr/bin/env python3
"""Build the SpatialLens Demo Day deck (.pptx) from the assets in
reports/demo_day/assets/ and the script in reports/demo_day/SLIDES_OUTLINE.md.

Run:  python scripts/build_demo_pptx.py
Out:  reports/demo_day/SpatialLens_DemoDay.pptx
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "reports" / "demo_day" / "assets"
OUT = ROOT / "reports" / "demo_day" / "SpatialLens_DemoDay.pptx"

# ---- palette -------------------------------------------------------------
INK = RGBColor(0x1A, 0x1A, 0x1A)      # near-black body text
NAVY = RGBColor(0x0B, 0x3D, 0x66)     # titles
ACCENT = RGBColor(0xE0, 0x6A, 0x00)   # orange — the numbers that matter
MUTE = RGBColor(0x5A, 0x5A, 0x5A)     # captions
PANEL = RGBColor(0xF2, 0xF4, 0xF7)    # light callout fill
FONT = "Calibri"

EMU_PER_IN = 914400


def fit(box_w_in, box_h_in, ar):
    """Return (w,h) in inches that fits an image of aspect ratio `ar`
    (w/h) inside the box, preserving aspect ratio."""
    if box_w_in / box_h_in > ar:
        h = box_h_in
        w = h * ar
    else:
        w = box_w_in
        h = w / ar
    return w, h


def img_ar(name):
    w, h = Image.open(ASSETS / name).size
    return w / h


def place_image(slide, name, left, top, box_w, box_h, align="center", valign="middle"):
    ar = img_ar(name)
    w, h = fit(box_w, box_h, ar)
    if align == "center":
        x = left + (box_w - w) / 2
    elif align == "right":
        x = left + (box_w - w)
    else:
        x = left
    if valign == "middle":
        y = top + (box_h - h) / 2
    elif valign == "bottom":
        y = top + (box_h - h)
    else:
        y = top
    return slide.shapes.add_picture(str(ASSETS / name), Inches(x), Inches(y),
                                    Inches(w), Inches(h))


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_run(run, text, size, color=INK, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def title_bar(slide, text, kicker=None):
    tb, tf = textbox(slide, 0.55, 0.30, 12.2, 1.0)
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, 30, NAVY, bold=True)
    if kicker:
        tb2, tf2 = textbox(slide, 0.57, 1.02, 12.0, 0.4)
        pp = tf2.paragraphs[0]
        set_run(pp.add_run(), kicker, 13, MUTE, italic=True)


def bullets(slide, items, left, top, width, height, size=16, gap=6):
    """items: list of (text, level, bold, color)"""
    tb, tf = textbox(slide, left, top, width, height)
    for i, it in enumerate(items):
        text, level, bold, color = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(gap)
        bullet = ("   • " if level == 0 else "        – ")
        set_run(p.add_run(), bullet + text, size, color, bold=bold)
    return tb


def callout(slide, left, top, width, height, text, size=22, fill=PANEL, color=ACCENT):
    box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text, size, color, bold=True)
    return box


def caption(slide, left, top, width, text, size=11, align=PP_ALIGN.CENTER):
    tb, tf = textbox(slide, left, top, width, 0.35)
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size, MUTE)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# -------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---- Slide 1: Title ------------------------------------------------------
s = prs.slides.add_slide(BLANK)
tb, tf = textbox(s, 0.7, 1.7, 7.6, 3.5)
p = tf.paragraphs[0]
set_run(p.add_run(), "SpatialLens Assist", 46, NAVY, bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(10)
set_run(p2.add_run(), "Motion-Aware Hazard Detection for Campus Navigation",
        24, INK, bold=True)
p3 = tf.add_paragraph()
set_run(p3.add_run(), "from Low-FPS Egocentric Video", 24, INK, bold=True)
p4 = tf.add_paragraph(); p4.space_before = Pt(22)
set_run(p4.add_run(), "Diego Sanchez  ·  Stanford CS131", 17, MUTE)
place_image(s, "slide1_title_thumb.jpg", 8.7, 0.7, 4.1, 6.1)
notes(s, "Hi, I'm Diego. SpatialLens Assist turns a cheap phone video of a campus "
         "walkway into spoken hazard alerts for blind and low-vision pedestrians — "
         "telling them not just WHAT is nearby, but whether it's COMING AT THEM. (15s)")


# ---- Slide 2: Problem ----------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_bar(s, "Detection tells you what, not what it’s doing")
bullets(s, [
    ("A parked bicycle and one bearing down on you are both just “bicycle, 0.86”.", 0, False, INK),
    ("A sighted person reads intent from context — a blind / low-vision", 0, False, INK),
    ("pedestrian can’t. The safety signal is motion, not the label.", 1, False, INK),
    ("Our task: per-track motion classification from a ~10 s clip @ 2 fps", 0, True, NAVY),
    ("approaching · crossing_LTR · crossing_RTL", 1, False, INK),
    ("moving_away · static · uncertain", 1, False, INK),
], 0.55, 1.55, 7.1, 4.8, size=18, gap=10)
place_image(s, "slide2_hazard_overlay.jpg", 8.0, 1.5, 4.8, 5.4)
caption(s, 8.0, 6.85, 4.8, "Hazard overlay on a campus walkway clip")
notes(s, "Here's the problem. A parked bike and a bike speeding toward you are detected "
         "identically — bicycle, point eight-six. A sighted person fills in the intent; "
         "the user we care about can't. So the safety-critical signal isn't the label, it's "
         "the motion. We frame it as per-track classification: given a ten-second clip at just "
         "two frames a second, decide whether each object is approaching, crossing left or "
         "right, moving away, static — or explicitly uncertain. (40s)")


# ---- Slide 3: Methodology ------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_bar(s, "A transparent, CPU-only pipeline")
place_image(s, "slide3_pipeline.png", 0.55, 1.35, 12.2, 2.5)
caption(s, 0.55, 3.85, 12.2,
        "2 fps frames  →  YOLOv8n  →  tracker + appearance re-ID  →  "
        "ego-motion-compensated motion features  →  rule cascade  →  alert + evidence string")
bullets(s, [
    ("Explainable — every prediction carries a natural-language evidence string.", 0, False, INK),
    ("Calibrated abstention — “uncertain” is first-class; short tracks don’t bluff.", 0, False, INK),
    ("CPU-only — the full pipeline runs in ~4 min on a GPU-less Mac.", 0, False, INK),
    ("Key trick: ego-motion compensation cancels camera pan, so static objects stay static.", 0, True, NAVY),
], 0.55, 4.35, 12.2, 2.7, size=17, gap=9)
notes(s, "The pipeline has six stages. Off-the-shelf YOLOv8n detects; a lightweight IoU-plus-"
         "centroid tracker links objects, with an appearance re-ID pass to stitch fragments. "
         "Then motion features — and the key trick is ego-motion compensation: we subtract "
         "the camera's own panning so a parked bench doesn't look like it's charging at you. "
         "Finally an interpretable rule cascade — no black box — emits the label plus a "
         "plain-English evidence string. Three principles: everything is explainable, it abstains "
         "rather than guesses, and it runs on CPU in about four minutes. (45s)")


# ---- Slide 4: Demo stage strip -------------------------------------------
s = prs.slides.add_slide(BLANK)
title_bar(s, "From pixels to a spoken alert")
labels = ["1 · Raw frame (2 fps)", "2 · YOLOv8n detections",
          "3 · Tracking + re-ID", "4 · Hazard overlay"]
files = ["slide4a_raw.jpg", "slide4b_detection.jpg", "slide4c_tracking.jpg", "slide4d_hazard.jpg"]
strip_top, strip_h = 1.45, 3.95
cell_w = 3.05
gap = 0.18
total = 4 * cell_w + 3 * gap
x0 = (13.333 - total) / 2
for i, (f, lab) in enumerate(zip(files, labels)):
    cx = x0 + i * (cell_w + gap)
    place_image(s, f, cx, strip_top, cell_w, strip_h, valign="top")
    caption(s, cx, strip_top + strip_h + 0.05, cell_w, lab, size=12)
callout(s, 2.4, 5.95, 8.5, 0.95,
        "\U0001F50A  “Bicycle approaching directly ahead.”   [high]", size=24)
caption(s, 2.4, 6.95, 8.5,
        "evidence string: the cyclist’s bounding box grew ~22×; centroid moved toward image center")
notes(s, "Here's a real run, left to right. Raw frame. YOLO's detections. The tracker links "
         "the cyclist across frames — that's the yellow trail. And the motion stage fires: "
         "the bounding box grew almost twenty-two-fold and the centroid moved toward the center, "
         "so the system says, out loud, 'bicycle approaching directly ahead,' high confidence — "
         "and it shows you exactly why it said that. (45s)")


# ---- Slide 5: Results ----------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_bar(s, "Accurate where it counts")
tb, tf = textbox(s, 0.55, 1.55, 6.4, 4.6)
def big(tf, first, num, rest):
    p = tf.add_paragraph(); p.space_after = Pt(14)
    set_run(p.add_run(), first, 19, INK)
    set_run(p.add_run(), num, 30, ACCENT, bold=True)
    if rest:
        set_run(p.add_run(), rest, 19, INK)
# clear default empty para
big(tf, "Overall accuracy:  ", "82.5%", "")
big(tf, "Decidable (≥3-frame tracks):  ", "93.3%", "")
big(tf, "approaching  P / R:  ", "0.95 / 0.88", "")
p = tf.add_paragraph(); p.space_before = Pt(6)
set_run(p.add_run(), "11 videos · 61 hand-labelled tracks", 14, MUTE, italic=True)
place_image(s, "slide5_confusion_matrix.png", 6.9, 1.5, 6.0, 4.5)
caption(s, 6.9, 6.0, 6.0, "Confusion matrix (57 evaluated tracks)")
callout(s, 0.55, 6.35, 6.1, 0.95,
        "Zero approaching ↔ moving_away confusion", size=18)
notes(s, "On 11 controlled walkway videos and 61 hand-labeled tracks: 82.5% overall, and "
         "93.3% on tracks long enough to decide — three frames or more. The safety-critical "
         "approaching class hits 0.95 precision, 0.88 recall. And the number I care most about: "
         "zero confusion between approaching and moving-away. When unsure it abstains into "
         "uncertain — it never told a user something was leaving when it was actually coming "
         "at them. For an alert system, that's the right asymmetry. (40s)")


# ---- Slide 6: Analysis ---------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_bar(s, "Continuity beats classifier complexity")
place_image(s, "slide6_ablation_study.png", 0.55, 1.6, 6.6, 4.9, valign="top")
caption(s, 0.55, 5.0, 6.6, "Accuracy drop when each component is removed")
bullets(s, [
    ("Biggest gains = keeping tracks intact, not fancier rules:", 0, True, NAVY),
    ("appearance re-ID:  +6.2 pts  (largest single gain)", 1, False, INK),
    ("two short-track salvage rules:  +3.5 pts each", 1, False, INK),
    ("Negative result: ByteTrack → 32.8% on 2 fps offline footage —", 0, True, NAVY),
    ("its online motion-continuity assumption breaks under sparse sampling.", 1, False, INK),
    ("Abstention = the “fragmentation tax made auditable”:", 0, False, INK),
    ("23 of 24 abstentions are short tracks, not classifier confusion.", 1, False, INK),
], 7.4, 1.6, 5.5, 5.2, size=16, gap=9)
notes(s, "The ablations taught me the real lesson. The biggest accuracy gains didn't come from "
         "smarter rules — they came from keeping tracks from fragmenting. Appearance re-ID "
         "alone is worth six points. And a useful negative result: ByteTrack, a state-of-the-art "
         "online tracker, drops to 33% here, because at two frames a second its motion-continuity "
         "assumption just doesn't hold. The system's uncertainty is almost entirely short broken "
         "tracks — so the abstentions are an honest, auditable tax, not confusion. (40s)")


# ---- Slide 7: Conclusion -------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_bar(s, "Takeaways & what’s next")
bullets(s, [
    ("Simple, interpretable motion cues recover useful hazard labels from cheap", 0, True, NAVY),
    ("low-fps video — and continuity matters more than classifier sophistication.", 0, True, NAVY),
    ("Next:", 0, False, INK),
    ("monocular depth (MiDaS)  ·  learned re-ID embedding", 1, False, INK),
    ("denser sampling (5–10 fps)  ·  user study with blind / low-vision walkers", 1, False, INK),
], 0.55, 1.8, 12.0, 3.6, size=20, gap=12)
tb, tf = textbox(s, 0.55, 5.7, 12.0, 1.0)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Thank you!", 30, NAVY, bold=True)
notes(s, "So: cheap, low-frame-rate video plus interpretable motion cues is enough to give "
         "useful, explainable hazard alerts. Next steps are adding depth, a learned re-ID "
         "embedding, and a study with actual blind and low-vision walkers. Thank you. (15s)")


prs.save(str(OUT))
print(f"wrote {OUT.relative_to(ROOT)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
